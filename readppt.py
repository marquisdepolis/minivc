import pptx
from pptx import Presentation
import os
import tempfile
from pptx.enum.shapes import MSO_SHAPE_TYPE

def read_text_from_shape(shape):
    text_runs = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
    return text_runs

def read_text_from_table(table):
    text_runs = []
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs

def process_grouped_shape(group_shape):
    text_runs = []
    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            text_runs.extend(process_grouped_shape(shape))
        elif shape.shape_type in [MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER]:
            text_runs.extend(read_text_from_shape(shape))
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            text_runs.extend(read_text_from_table(shape.table))
    return text_runs

def read_ppt(file_content):
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        temp_file.write(file_content)
        temp_file.flush()

        presentation = Presentation(temp_file.name)
        text_runs = []

        for idx, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type in [MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER]:
                    text_runs.extend(read_text_from_shape(shape))
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    text_runs.extend(read_text_from_table(shape.table))
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    text_runs.extend(process_grouped_shape(shape))

        text = " ".join(text_runs)
        os.unlink(temp_file.name)

    return text